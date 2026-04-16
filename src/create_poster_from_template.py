from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Suggested-Template-for-POSTER(1).pptx"
OUTPUT = ROOT / "Capstone_Poster_TemplateBased_Wenxuan_Li.pptx"

ROC = ROOT / "outputs/figures/resnet50_transfer_full_v1/test_roc_curve.png"
GRADCAM = ROOT / "outputs/figures/resnet50_transfer_full_v1/gradcam/gradcam_overview.png"

BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(70, 70, 70)
LIGHT = RGBColor(244, 246, 248)


def fill_text(shape, lines, font_size, bold=False, align=PP_ALIGN.LEFT, color=BLACK, spacing=1.05):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, lines, font_size, bold=False, align=PP_ALIGN.LEFT, color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def style_table(table):
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            cell.fill.background()
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(20 if r == 0 else 18)
                    run.font.bold = r == 0
                    run.font.color.rgb = BLACK


def add_workflow_block(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = GRAY
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = BLACK


def add_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY
    arrow.line.fill.background()


def main():
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    shapes = list(slide.shapes)
    title = shapes[0]
    author = shapes[1]
    dept = shapes[2]
    intro_body = shapes[3]
    results_body = shapes[4]
    conclusions_body = shapes[5]
    methods_body = shapes[6]
    refs_body = shapes[7]

    fill_text(title, ["Deep Learning for Medical Image Processing"], 42, bold=True, align=PP_ALIGN.CENTER)
    fill_text(author, ["Wenxuan Li, Yue Zhao"], 28, bold=True, align=PP_ALIGN.CENTER)
    fill_text(
        dept,
        ["Department of Computer Science, Wenzhou-Kean University, liwenxu@kean.edu"],
        20,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    fill_text(
        intro_body,
        [
            "• Chest X-ray screening is widely used, but manual interpretation is time-consuming and reader-dependent.",
            "• This project studies whether transfer learning and Grad-CAM improve binary chest X-ray classification.",
            "• Main question: which model and training strategy gives the best result?",
            "• Can Grad-CAM explain representative correct and incorrect predictions?",
        ],
        24,
        spacing=1.1,
    )

    fill_text(
        methods_body,
        [
            "• Dataset: NIH ChestXray14, 112,120 images.",
            "• Task: Normal (No Finding) vs. Abnormal (any pathology).",
            "• Models: ResNet50 and DenseNet121, each in scratch and transfer-learning settings.",
            "• Metrics: accuracy, F1, AUC, sensitivity, specificity, ROC, and Grad-CAM.",
        ],
        22,
        spacing=1.08,
    )

    m_left, m_top, m_width, _ = methods_body.left, methods_body.top, methods_body.width, methods_body.height
    wf_y = m_top + Inches(3.7)
    block_w = Inches(2.15)
    block_h = Inches(1.05)
    gap = Inches(0.55)
    b1_x = m_left + Inches(0.4)
    b2_x = b1_x + block_w + gap
    b3_x = b2_x + block_w + gap
    b4_x = b3_x + block_w + gap
    add_workflow_block(slide, b1_x, wf_y, block_w, block_h, "NIH Dataset")
    add_arrow(slide, b1_x + block_w + Inches(0.08), wf_y + Inches(0.26), Inches(0.35), Inches(0.28))
    add_workflow_block(slide, b2_x, wf_y, block_w, block_h, "Preprocess\n224 x 224")
    add_arrow(slide, b2_x + block_w + Inches(0.08), wf_y + Inches(0.26), Inches(0.35), Inches(0.28))
    add_workflow_block(slide, b3_x, wf_y, block_w, block_h, "4 Training\nRuns")
    add_arrow(slide, b3_x + block_w + Inches(0.08), wf_y + Inches(0.26), Inches(0.35), Inches(0.28))
    add_workflow_block(slide, b4_x, wf_y, block_w, block_h, "Metrics +\nGrad-CAM")
    add_textbox(
        slide,
        m_left + Inches(0.45),
        wf_y + Inches(1.35),
        m_width - Inches(0.9),
        Inches(1.4),
        [
            "Training setup: batch size 16, Adam, BCEWithLogitsLoss, max 20 epochs,",
            "early stopping on validation AUC.",
        ],
        20,
        align=PP_ALIGN.CENTER,
        color=GRAY,
    )

    fill_text(
        conclusions_body,
        [
            "• DenseNet121 scratch achieved the highest test AUC (0.7457).",
            "• ResNet50 transfer gave the best overall balance: accuracy 0.6955, F1 0.6547, AUC 0.7454.",
            "• Transfer learning helped ResNet50 clearly, but not uniformly DenseNet121.",
            "• Grad-CAM was useful for qualitative error analysis.",
        ],
        24,
        spacing=1.1,
    )

    fill_text(
        refs_body,
        [
            "Wang et al., CVPR 2017, ChestX-ray8/ChestXray14.",
            "He et al., CVPR 2016, ResNet.",
            "Huang et al., CVPR 2017, DenseNet.",
            "Selvaraju et al., ICCV 2017, Grad-CAM.",
        ],
        18,
        spacing=1.0,
    )

    results_body.text = ""
    left, top, width, _ = results_body.left, results_body.top, results_body.width, results_body.height

    add_textbox(
        slide,
        left + Inches(0.15),
        top + Inches(0.12),
        width - Inches(0.3),
        Inches(0.95),
        [
            "Best AUC: DenseNet121 scratch (0.7457)   |   Best balance: ResNet50 transfer (Acc 0.6955, F1 0.6547)",
        ],
        22,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    table_shape = slide.shapes.add_table(
        5,
        5,
        left + Inches(0.4),
        top + Inches(1.05),
        width - Inches(0.8),
        Inches(2.55),
    )
    table = table_shape.table
    table.columns[0].width = Inches(2.1)
    table.columns[1].width = Inches(2.1)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(1.7)
    table.columns[4].width = Inches(1.7)
    headers = ["Model", "Training", "Accuracy", "F1", "AUC"]
    rows = [
        ["ResNet50", "Scratch", "0.673", "0.612", "0.723"],
        ["ResNet50", "Transfer", "0.695", "0.655", "0.745"],
        ["DenseNet121", "Scratch", "0.691", "0.639", "0.746"],
        ["DenseNet121", "Transfer", "0.671", "0.565", "0.742"],
    ]
    for c, text in enumerate(headers):
        table.cell(0, c).text = text
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    style_table(table)

    image_top = top + Inches(4.0)
    if ROC.exists():
        slide.shapes.add_picture(str(ROC), left + Inches(0.4), image_top, width=Inches(4.95))
    if GRADCAM.exists():
        slide.shapes.add_picture(str(GRADCAM), left + Inches(5.75), image_top, width=Inches(4.95))

    add_textbox(
        slide,
        left + Inches(0.3),
        top + Inches(12.3),
        width - Inches(0.6),
        Inches(1.25),
        [
            "Left: ROC curve for ResNet50 transfer. Right: Grad-CAM examples for correct and error cases.",
            "Quantitative metrics and visual explanations together show meaningful tradeoffs across the four experiments.",
        ],
        18,
        align=PP_ALIGN.CENTER,
        color=GRAY,
    )

    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT}")


if __name__ == "__main__":
    main()
