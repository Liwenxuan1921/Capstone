from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "Capstone_Poster_Wenxuan_Li.pptx"
LOGO = ROOT / "template_preview.png"
ROC = ROOT / "outputs/figures/resnet50_transfer_full_v1/test_roc_curve.png"
CM = ROOT / "outputs/figures/resnet50_transfer_full_v1/test_confusion_matrix.png"
GRADCAM = ROOT / "outputs/figures/resnet50_transfer_full_v1/gradcam/gradcam_overview.png"


KEAN_BLUE = RGBColor(0, 61, 124)
LIGHT_BLUE = RGBColor(232, 241, 248)
MID_BLUE = RGBColor(210, 227, 241)
DARK = RGBColor(34, 34, 34)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, left, top, width, height, text="", font_size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, fill=None, line=None,
                margin=0.08, font_name="Calibri"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.adjustments[0] = 0.02
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.color.rgb = KEAN_BLUE
        shape.line.width = Pt(1.2)
    elif line is False:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_paragraphs(shape, paragraphs, font_size=16, color=DARK, bullet=True, level=0):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.bullet = bullet
        p.space_after = Pt(4)
        p.font.name = "Calibri"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color


def add_section(slide, left, top, width, height, title):
    body = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    body.adjustments[0] = 0.02
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = KEAN_BLUE
    body.line.width = Pt(1.6)

    header_h = Inches(0.65)
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, header_h
    )
    header.adjustments[0] = 0.02
    header.fill.solid()
    header.fill.fore_color.rgb = KEAN_BLUE
    header.line.fill.background()
    tf = header.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE

    return body


def add_image(slide, path, left, top, width=None, height=None):
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def build_table(slide, left, top, width, height):
    rows, cols = 5, 5
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    col_widths = [1.8, 1.8, 1.4, 1.3, 1.3]
    for idx, w in enumerate(col_widths):
        table.columns[idx].width = Inches(w)

    headers = ["Model", "Training", "Acc.", "F1", "AUC"]
    data = [
        ["ResNet50", "Scratch", "0.673", "0.612", "0.723"],
        ["ResNet50", "Transfer", "0.695", "0.655", "0.745"],
        ["DenseNet121", "Scratch", "0.691", "0.639", "0.746"],
        ["DenseNet121", "Transfer", "0.671", "0.565", "0.742"],
    ]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = KEAN_BLUE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Calibri"
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.color.rgb = WHITE

    for r_idx, row in enumerate(data, start=1):
        for c_idx, text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if r_idx % 2 == 1 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(13)
                    run.font.color.rgb = DARK
                    if row[0] == "DenseNet121" and row[1] == "Scratch" and c_idx == 4:
                        run.font.bold = True
                    if row[0] == "ResNet50" and row[1] == "Transfer" and c_idx in (2, 3):
                        run.font.bold = True


def main():
    prs = Presentation()
    prs.slide_width = Inches(56)
    prs.slide_height = Inches(36)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 250, 252)

    # Header band
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(4.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = KEAN_BLUE
    header.line.fill.background()

    if LOGO.exists():
        add_image(slide, LOGO, Inches(0.7), Inches(0.45), height=Inches(2.9))

    add_textbox(
        slide, Inches(4.2), Inches(0.45), Inches(48.8), Inches(1.7),
        "Deep Learning for Medical Image Processing:\nImproving Chest X-ray Diagnostic Classification\nwith Transfer Learning and Explainable Deep Learning",
        font_size=28, bold=True, color=WHITE, fill=None, line=False, margin=0.02
    )
    add_textbox(
        slide, Inches(4.25), Inches(2.35), Inches(48.5), Inches(0.7),
        "Wenxuan Li  |  Supervisor: Yue Zhao  |  Department of Computer Science and Technology, Wenzhou-Kean University",
        font_size=17, color=WHITE, fill=None, line=False, margin=0.02
    )
    add_textbox(
        slide, Inches(4.25), Inches(3.0), Inches(48.5), Inches(0.55),
        "Capstone Poster Summary",
        font_size=16, bold=True, color=RGBColor(219, 236, 255), fill=None, line=False, margin=0.02
    )

    col_gap = Inches(0.45)
    col_w = Inches(18.1)
    x1 = Inches(0.7)
    x2 = x1 + col_w + col_gap
    x3 = x2 + col_w + col_gap

    # Left column
    add_section(slide, x1, Inches(4.6), col_w, Inches(6.5), "Introduction")
    intro = slide.shapes.add_textbox(x1 + Inches(0.2), Inches(5.35), col_w - Inches(0.4), Inches(5.4))
    add_paragraphs(
        intro,
        [
            "Chest X-ray screening is widely used, but manual interpretation is time-consuming and subject to reader variability.",
            "This project tests whether transfer learning and Grad-CAM can improve binary chest X-ray classification on NIH ChestXray14.",
            "Focus: compare ResNet50 and DenseNet121, then analyze errors visually.",
        ],
        font_size=16,
    )

    add_section(slide, x1, Inches(11.45), col_w, Inches(8.8), "Methods and Materials")
    methods = slide.shapes.add_textbox(x1 + Inches(0.2), Inches(12.2), col_w - Inches(0.4), Inches(7.8))
    add_paragraphs(
        methods,
        [
            "Dataset: NIH ChestXray14, 112,120 images.",
            "Task: Normal (No Finding) vs. Abnormal (any pathology label).",
            "Models: ResNet50 and DenseNet121, each in scratch and transfer-learning settings.",
            "Setup: 224×224 input, batch size 16, Adam, BCEWithLogitsLoss, max 20 epochs, early stopping on validation AUC.",
            "Outputs: accuracy, F1, AUC, ROC, confusion matrix, and Grad-CAM.",
        ],
        font_size=15,
    )

    add_section(slide, x1, Inches(20.6), col_w, Inches(8.0), "Key Takeaways")
    takeaways = slide.shapes.add_textbox(x1 + Inches(0.2), Inches(21.35), col_w - Inches(0.4), Inches(6.9))
    add_paragraphs(
        takeaways,
        [
            "DenseNet121 from scratch achieved the highest test AUC (0.7457).",
            "ResNet50 with transfer learning provided the strongest overall balance of accuracy (0.6955), recall (0.6272), and F1-score (0.6547).",
            "Transfer learning clearly benefited ResNet50, but not uniformly DenseNet121.",
            "Grad-CAM highlighted plausible thoracic regions, while error cases exposed remaining limitations.",
        ],
        font_size=16,
    )

    # Middle column
    add_section(slide, x2, Inches(4.6), col_w, Inches(10.2), "Results")
    build_table(slide, x2 + Inches(0.35), Inches(5.45), Inches(8.2), Inches(3.0))
    add_textbox(
        slide, x2 + Inches(8.9), Inches(5.25), Inches(8.7), Inches(3.2),
        "Best AUC:\nDenseNet121 Scratch\n0.7457\n\nBest balance:\nResNet50 Transfer\nAcc 0.6955 | F1 0.6547 | AUC 0.7454",
        font_size=17, bold=True, color=KEAN_BLUE, fill=LIGHT_BLUE, line=KEAN_BLUE, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, x2 + Inches(0.35), Inches(8.95), Inches(17.35), Inches(1.1),
        "DenseNet121 scratch achieved the highest AUC, while ResNet50 transfer produced the most balanced classification behavior.",
        font_size=16, fill=None, line=False, margin=0.02
    )

    add_section(slide, x2, Inches(15.1), col_w, Inches(13.5), "Model Performance Visualizations")
    if ROC.exists():
        add_image(slide, ROC, x2 + Inches(0.35), Inches(16.05), width=Inches(8.2))
    if CM.exists():
        add_image(slide, CM, x2 + Inches(9.0), Inches(16.05), width=Inches(8.0))
    add_textbox(
        slide, x2 + Inches(0.35), Inches(24.35), Inches(17.2), Inches(2.0),
        "Left: ROC curve for ResNet50 with transfer learning.\nRight: Confusion matrix on the test set.",
        font_size=15, fill=None, line=False, margin=0.02
    )

    # Right column
    add_section(slide, x3, Inches(4.6), col_w, Inches(13.2), "Explainability (Grad-CAM)")
    if GRADCAM.exists():
        add_image(slide, GRADCAM, x3 + Inches(0.35), Inches(5.45), width=Inches(17.2))
    add_textbox(
        slide, x3 + Inches(0.35), Inches(14.95), Inches(17.2), Inches(1.9),
        "Representative Grad-CAM outputs for correct and error cases. Heatmaps often focus on thoracic regions, but they must be interpreted alongside quantitative results.",
        font_size=15, fill=None, line=False, margin=0.02
    )

    add_section(slide, x3, Inches(18.15), col_w, Inches(6.3), "Conclusions")
    conclusions = slide.shapes.add_textbox(x3 + Inches(0.2), Inches(18.95), col_w - Inches(0.4), Inches(5.1))
    add_paragraphs(
        conclusions,
        [
            "Deep learning supports binary chest X-ray classification effectively in a reproducible capstone workflow.",
            "The best model depends on the metric of interest: highest AUC vs. best overall balance.",
            "Grad-CAM is useful for qualitative inspection and failure-case analysis.",
        ],
        font_size=16,
    )

    add_section(slide, x3, Inches(24.75), col_w, Inches(3.85), "References")
    refs = slide.shapes.add_textbox(x3 + Inches(0.2), Inches(25.45), col_w - Inches(0.4), Inches(2.95))
    refs.text_frame.word_wrap = True
    ref_lines = [
        "Wang et al., CVPR 2017, ChestX-ray8/ChestXray14.",
        "He et al., CVPR 2016, ResNet.",
        "Huang et al., CVPR 2017, DenseNet.",
        "Selvaraju et al., ICCV 2017, Grad-CAM.",
        "Additional 2024-2026 chest X-ray and XAI studies are cited in the thesis.",
    ]
    add_paragraphs(refs, ref_lines, font_size=12.5, bullet=False)

    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT}")


if __name__ == "__main__":
    main()
